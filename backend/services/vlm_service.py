# backend/services/vlm_service.py

import os
import io
import re
import unicodedata

import fitz
from pptx import Presentation
from PyPDF2 import PdfReader
from docx import Document
from PIL import Image

import torch
from transformers import AutoProcessor, LlavaOnevisionForConditionalGeneration


# ========================================
# 파일명/텍스트 정리 함수
# ========================================
def sanitize_filename(name: str) -> str:
    name = unicodedata.normalize("NFKC", name)
    name = re.sub(r'[\\/:*?"<>|]', "_", name)
    return name.strip(" .")


def clean_text(text):
    text = re.sub(r'[^\x09\x0A\x0D\x20-\uFFFF]', '', text)
    text = re.sub(r'[\uF000-\uF2FF]', '', text)
    text = re.sub(r'[•●▪★☆◆◇○◉◎∎]', '', text)
    return text.replace("�", "")


# ========================================
# 1) Lazy Loading: 요청 시 로드
# ========================================
device = "cuda" if torch.cuda.is_available() else "cpu"
model_name = "NCSOFT/VARCO-VISION-2.0-1.7B"

processor = None
model = None


def load_vlm_model():
    """모델을 최초 1회만 로드하는 함수"""
    global processor, model

    if processor is not None and model is not None:
        return  # 이미 로드됨

    print("\n[VLM] Loading model... (최초 1회)\n")

    processor = AutoProcessor.from_pretrained(model_name)

    model = LlavaOnevisionForConditionalGeneration.from_pretrained(
        model_name,
        dtype=torch.float16 if device == "cuda" else torch.float32,
    ).to(device)

    model.config.pad_token_id = model.config.eos_token_id

    print("\n[VLM] Model loaded successfully!\n")


# ========================================
# 2) 텍스트 블록 추출
# ========================================
def extract_pdf_text_blocks(path: str) -> list[str]:
    reader = PdfReader(path)
    blocks = []
    for idx, page in enumerate(reader.pages):
        page_text = page.extract_text() or ""
        blocks.append(f"=== PAGE {idx+1} ===\n{page_text}\n")
    return blocks


def _extract_shapes(shape):
    texts = []
    if getattr(shape, "has_text_frame", False):
        texts.append(shape.text)
    if hasattr(shape, "shapes"):
        for shp in shape.shapes:
            texts.extend(_extract_shapes(shp))
    return texts


def extract_pptx_text_blocks(path: str) -> list[str]:
    pres = Presentation(path)
    blocks = []

    for idx, slide in enumerate(pres.slides):
        slide_texts = []
        for shape in slide.shapes:
            slide_texts.extend(_extract_shapes(shape))

        joined_text = "\n".join(slide_texts)
        block = f"=== SLIDE {idx+1} ===\n{joined_text}\n"
        blocks.append(block)

    return blocks


# ========================================
# 3) 이미지 추출
# ========================================
def extract_images_pdf(path: str):
    doc = fitz.open(path)
    images = []

    for page_index in range(len(doc)):
        page = doc[page_index]
        for img_i, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base = doc.extract_image(xref)
            pil_img = Image.open(io.BytesIO(base["image"])).convert("RGB")
            images.append((page_index + 1, img_i + 1, pil_img))

    return images


def extract_images_pptx(path: str):
    pres = Presentation(path)
    images = []

    for i, slide in enumerate(pres.slides):
        img_index = 1
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                pil_img = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                images.append((i + 1, img_index, pil_img))
                img_index += 1

    return images


# ========================================
# 4) 캡션 생성
# ========================================
def clean_caption(text: str) -> str:
    remove = ["user", "assistant", "<|im_start|>", "<|im_end|>"]
    for r in remove:
        text = text.replace(r, "")
    return text.strip()


def generate_caption(image_pil: Image.Image) -> str:
    load_vlm_model()  # 요청 시점에 모델 로드

    prompt = """<|im_start|>user
                <image>
                이 이미지를 한국어로 자세히 설명해줘.
                <|im_end|>
                <|im_start|>assistant
            """

    inputs = processor(
        images=image_pil,
        text=prompt,
        return_tensors="pt",
    ).to(device)

    output = model.generate(**inputs, max_new_tokens=120)
    raw = processor.batch_decode(output, skip_special_tokens=True)[0]
    return clean_caption(raw)


# ========================================
# 5) 텍스트 + 이미지 캡션 합치기
# ========================================
def build_text_with_captions(input_path: str):
    lower = input_path.lower()

    if lower.endswith(".pdf"):
        text_blocks = extract_pdf_text_blocks(input_path)
        images = extract_images_pdf(input_path)
        prefix = "PAGE"
    elif lower.endswith(".pptx"):
        text_blocks = extract_pptx_text_blocks(input_path)
        images = extract_images_pptx(input_path)
        prefix = "SLIDE"
    else:
        raise ValueError("pdf 또는 pptx만 지원합니다.")

    img_dict = {}
    for page_num, img_index, pil_img in images:
        img_dict.setdefault(page_num, []).append((img_index, pil_img))

    final_lines = []

    for idx, raw_block in enumerate(text_blocks):
        page_num = idx + 1
        final_lines.append(clean_text(raw_block))

        if page_num in img_dict:
            for img_index, pil_img in img_dict[page_num]:
                caption = generate_caption(pil_img)
                final_lines.append(
                    f"=== {prefix} {page_num} IMAGE {img_index} ===\n{caption}\n"
                )

    return final_lines, prefix


# ========================================
# 6) 최종 파일 저장
# ========================================
def convert_to_text_and_docx_with_vlm(input_path: str, output_dir="static", desired_name=None):
    os.makedirs(output_dir, exist_ok=True)

    if desired_name:
        base_name = sanitize_filename(desired_name)
    else:
        base_name = os.path.splitext(os.path.basename(input_path))[0]

    final_lines, _ = build_text_with_captions(input_path)

    txt_path = os.path.join(output_dir, base_name + "_vlm.txt")
    docx_path = os.path.join(output_dir, base_name + "_vlm.docx")

    with open(txt_path, "w", encoding="utf-8") as f:
        f.write("\n".join(final_lines))

    doc = Document()
    for line in final_lines:
        doc.add_paragraph(line)
    doc.save(docx_path)

    return docx_path
