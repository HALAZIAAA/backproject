# backend/services/convert_service.py

import re, os, unicodedata
from datetime import datetime
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document

def sanitize_filename(name: str) -> str:
    # 확장자 제거된 기본 이름만 받는다고 가정
    name = unicodedata.normalize("NFKC", name)
    # 윈도우 금지문자 제거: \ / : * ? " < > |
    name = re.sub(r'[\\/:*?"<>|]', "_", name)
    # 앞뒤 공백/점 제거
    return name.strip(" .")

def extract_pdf_text(path):
    reader = PdfReader(path)
    out = []
    for idx, page in enumerate(reader.pages):
        page_text = page.extract_text() or ""
        out.append(f"=== PAGE {idx+1} ===\n{page_text}\n")
    return "\n".join(out)

def extract_shapes(shape):
    texts = []
    if shape.has_text_frame:
        texts.append(shape.text)
    if hasattr(shape, "shapes"):
        for shp in shape.shapes:
            texts.extend(extract_shapes(shp))
    return texts

def extract_pptx_text(path):
    pres = Presentation(path)
    out = []
    for idx, slide in enumerate(pres.slides):
        slide_texts = []
        for shape in slide.shapes:
            slide_texts.extend(extract_shapes(shape))
        joined = "\n".join(slide_texts)
        out.append(f"=== SLIDE {idx+1} ===\n{joined}\n")
    return "\n".join(out)

def clean_text(text):
    text = re.sub(r'[^\x09\x0A\x0D\x20-\uFFFF]', '', text)
    text = re.sub(r'[\uF000-\uF2FF]', '', text)
    text = re.sub(r'[•●▪★☆◆◇○◉◎∎]', '', text)
    text = text.replace("�", "")
    return text

def convert_to_text_and_docx(input_path, output_dir="static", desired_name: str | None = None):
    """
    PDF/PPTX → txt/docx 변환.
    desired_name: 확장자 없는 '원래 파일명'을 넘기면 그 이름으로 저장.
    """
    os.makedirs(output_dir, exist_ok=True)

    # 저장 파일명 결정
    if desired_name:
        base_name = sanitize_filename(desired_name)
    else:
        base_name = os.path.splitext(os.path.basename(input_path))[0]

    # (선택) 중복 방지용 타임스탬프 붙이려면 아래 주석 해제
    # ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    # base_name = f"{base_name}_{ts}"

    # 변환
    lower = input_path.lower()
    if lower.endswith(".pdf"):
        text = extract_pdf_text(input_path)
    elif lower.endswith(".pptx"):
        text = extract_pptx_text(input_path)
    else:
        raise ValueError("pdf 또는 pptx만 지원합니다.")

    text = clean_text(text)

    txt_path  = os.path.join(output_dir, base_name + ".txt")
    docx_path = os.path.join(output_dir, base_name + ".docx")

    with open(txt_path, "w", encoding="utf-8") as f:
        f.write(text)

    doc = Document()
    for line in text.split("\n"):
        doc.add_paragraph(line)
    doc.save(docx_path)

    return txt_path, docx_path
